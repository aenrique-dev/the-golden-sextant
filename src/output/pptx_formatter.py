"""
pptx_formatter.py — Branded PowerPoint output for The Partnership Protocol

Generates a .pptx digest file using Obsidian brand colors and logo:
  - Slide 1: Cover (title, date, exec summary, run stats)
  - Slides 2–N: One slide per DigestEntry (signal + action + outreach draft)

Brand palette (from Obsidian Mini Brand Guide):
  Primary:   #0B173B (navy), #002594 (deep blue), #FFFFFF
  Secondary: #0431B6, #5565E2, #61CEF4 (cyan), #FF9D03 (orange)
  Neutral:   #131313, #4E4E4E, #898989, #C4C4C4

Usage:
    from output.pptx_formatter import to_pptx
    path = to_pptx(digest, "data/digests/digest-2026-03-13-abc123.pptx")
"""

from pathlib import Path
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from models.schemas import Digest, DigestEntry, Priority


# ── Brand colors ────────────────────────────────────────────────────────────────

BG_NAVY    = RGBColor(0x0B, 0x17, 0x3B)   # #0B173B  primary dark background
CARD_BG    = RGBColor(0x0D, 0x1F, 0x50)   # #0D1F50  slightly lighter navy for cards
BLUE_DEEP  = RGBColor(0x00, 0x25, 0x94)   # #002594  deep blue
BLUE_BRITE = RGBColor(0x04, 0x31, 0xB6)   # #0431B6  bright blue (accent bar)
BLUE_PERI  = RGBColor(0x55, 0x65, 0xE2)   # #5565E2  periwinkle (labels, medium)
CYAN       = RGBColor(0x61, 0xCE, 0xF4)   # #61CEF4  cyan (company names, subtitles)
ORANGE     = RGBColor(0xFF, 0x9D, 0x03)   # #FF9D03  orange (high priority)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)   # #FFFFFF
MUTED      = RGBColor(0x89, 0x89, 0x89)   # #898989  muted / footer text
LIGHT_GRAY = RGBColor(0xC4, 0xC4, 0xC4)   # #C4C4C4  secondary body text

PRIORITY_COLORS = {
    Priority.HIGH:   ORANGE,
    Priority.MEDIUM: BLUE_PERI,
    Priority.LOW:    MUTED,
}

SIGNAL_TYPE_LABELS = {
    "competitor_move":       "Competitor Move",
    "partner_announcement":  "Partner Announcement",
    "tech_integration":      "Tech Integration",
    "marketplace_listing":   "Marketplace Listing",
    "joint_customer_story":  "Joint Customer Story",
    "webinar_event":         "Webinar / Event",
    "mssp_motion":           "MSSP Motion",
    "category_trend":        "Category Trend",
    "unknown":               "Signal",
}

# Logo — white version for dark slides, relative to this file's location
LOGO_PATH = Path(__file__).parent.parent.parent / "images" / "Obsidian Logo - White@2x.png"

# Slide dimensions: 16:9
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── Layout helpers ──────────────────────────────────────────────────────────────

def _set_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_rect(slide, left, top, width, height, fill_color: RGBColor, line=False):
    """Add a filled rectangle with no border (unless line=True)."""
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line:
        shape.line.color.rgb = fill_color
    else:
        shape.line.fill.background()
    return shape


def _add_text(slide, text, left, top, width, height,
              size=Pt(12), bold=False, color=WHITE,
              align=PP_ALIGN.LEFT, wrap=True):
    """Add a textbox with a single run."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return box


def _add_logo(slide, left=Inches(0.3), top=Inches(0.18), width=Inches(1.7)):
    if LOGO_PATH.exists():
        slide.shapes.add_picture(str(LOGO_PATH), left, top, width=width)


# ── Cover slide ─────────────────────────────────────────────────────────────────

def _build_cover(prs: Presentation, digest: Digest):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _set_bg(slide, BG_NAVY)

    # Left accent bar
    _add_rect(slide, Inches(0), Inches(0), Inches(0.07), SLIDE_H, BLUE_BRITE)

    # Logo — top left
    _add_logo(slide)

    # Title
    _add_text(
        slide, "The Partnership Protocol",
        Inches(0.5), Inches(1.6), Inches(12), Inches(1.1),
        size=Pt(46), bold=True, color=WHITE
    )

    # Date + run ID
    now = datetime.utcnow().strftime("%B %d, %Y")
    _add_text(
        slide, f"Daily Ecosystem Digest  ·  {now}  ·  Run {digest.run_id}",
        Inches(0.5), Inches(2.7), Inches(11), Inches(0.45),
        size=Pt(15), color=CYAN
    )

    # Exec summary
    _add_text(
        slide, digest.summary_narrative,
        Inches(0.5), Inches(3.35), Inches(11.8), Inches(1.3),
        size=Pt(13), color=LIGHT_GRAY, wrap=True
    )

    # Stats pills background
    _add_rect(slide, Inches(0.5), Inches(4.9), Inches(11.3), Inches(0.55), CARD_BG)
    stats = (
        f"  {digest.total_pages_scanned} pages scanned   ·   "
        f"{digest.total_signals_found} signals found   ·   "
        f"{len(digest.top_entries)} signals in digest"
    )
    _add_text(
        slide, stats,
        Inches(0.5), Inches(4.93), Inches(11.3), Inches(0.45),
        size=Pt(13), color=CYAN
    )

    # Bottom rule + footer
    _add_rect(slide, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.04), BLUE_DEEP)
    _add_text(
        slide, "Powered by Claude · Anthropic",
        Inches(0.5), Inches(7.1), Inches(6), Inches(0.35),
        size=Pt(9), color=MUTED
    )


# ── Signal slide ────────────────────────────────────────────────────────────────

def _build_signal_slide(prs: Presentation, entry: DigestEntry, idx: int, total: int):
    s = entry.signal
    action = entry.actions[0] if entry.actions else None

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, BG_NAVY)

    # Left accent bar
    _add_rect(slide, Inches(0), Inches(0), Inches(0.07), SLIDE_H, BLUE_BRITE)

    # Signal counter (top-left)
    _add_text(
        slide, f"Signal {idx} of {total}",
        Inches(0.2), Inches(0.08), Inches(3), Inches(0.3),
        size=Pt(8), color=MUTED
    )

    # Priority badge (top-right)
    p_color = PRIORITY_COLORS.get(s.priority, BLUE_PERI)
    badge_left = SLIDE_W - Inches(1.45)
    _add_rect(slide, badge_left, Inches(0.1), Inches(1.3), Inches(0.38), p_color)
    _add_text(
        slide, s.priority.upper(),
        badge_left, Inches(0.1), Inches(1.3), Inches(0.38),
        size=Pt(11), bold=True, color=WHITE, align=PP_ALIGN.CENTER
    )

    # Company + type (top row)
    type_label = SIGNAL_TYPE_LABELS.get(s.signal_type, "Signal")
    _add_text(
        slide, f"{s.company_name.upper()}  ·  {type_label}",
        Inches(0.5), Inches(0.15), Inches(10.5), Inches(0.38),
        size=Pt(11), bold=True, color=CYAN
    )

    # Signal title
    _add_text(
        slide, s.title,
        Inches(0.5), Inches(0.6), Inches(12.3), Inches(1.0),
        size=Pt(22), bold=True, color=WHITE, wrap=True
    )

    # Divider
    _add_rect(slide, Inches(0.5), Inches(1.68), Inches(12.3), Inches(0.03), BLUE_DEEP)

    # What happened
    _add_text(slide, "WHAT HAPPENED",
              Inches(0.5), Inches(1.8), Inches(4), Inches(0.28),
              size=Pt(8), bold=True, color=BLUE_PERI)
    _add_text(slide, s.summary,
              Inches(0.5), Inches(2.08), Inches(12.3), Inches(0.9),
              size=Pt(12), color=WHITE, wrap=True)

    # Why it matters
    _add_text(slide, "WHY IT MATTERS",
              Inches(0.5), Inches(3.08), Inches(4), Inches(0.28),
              size=Pt(8), bold=True, color=BLUE_PERI)
    _add_text(slide, s.why_it_matters,
              Inches(0.5), Inches(3.36), Inches(12.3), Inches(0.8),
              size=Pt(12), color=WHITE, wrap=True)

    # Action card
    card_top = Inches(4.28)
    _add_rect(slide, Inches(0.3), card_top, Inches(12.73), Inches(2.65), CARD_BG)

    if action:
        # Header row: "RECOMMENDED ACTION" label + action type chip
        _add_text(slide, "RECOMMENDED ACTION",
                  Inches(0.5), card_top + Inches(0.1), Inches(5), Inches(0.28),
                  size=Pt(8), bold=True, color=BLUE_PERI)

        action_label = action.action_type.replace("_", " ").title()
        chip_left = Inches(6.8)
        _add_rect(slide, chip_left, card_top + Inches(0.08),
                  Inches(2.8), Inches(0.3), BLUE_BRITE)
        _add_text(slide, action_label,
                  chip_left, card_top + Inches(0.08), Inches(2.8), Inches(0.3),
                  size=Pt(8), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        # Action description
        _add_text(slide, action.description,
                  Inches(0.5), card_top + Inches(0.45), Inches(12.3), Inches(0.65),
                  size=Pt(12), bold=True, color=WHITE, wrap=True)

        # Owner + partners
        owner_line = f"Owner: {action.suggested_owner}"
        if action.suggested_partners:
            owner_line += f"   ·   Activate: {', '.join(action.suggested_partners)}"
        _add_text(slide, owner_line,
                  Inches(0.5), card_top + Inches(1.12), Inches(12.3), Inches(0.35),
                  size=Pt(10), color=CYAN)

        # Outreach draft (if present)
        if action.outreach_draft:
            draft = f'"{action.outreach_draft}"'
            _add_text(slide, draft,
                      Inches(0.5), card_top + Inches(1.5), Inches(12.3), Inches(1.0),
                      size=Pt(10), color=LIGHT_GRAY, wrap=True)

    # Footer: priority rationale + source
    footer = s.source_url
    if s.priority_rationale:
        footer = f"{s.priority_rationale}   ·   {s.source_url}"
    _add_text(slide, footer,
              Inches(0.5), Inches(7.15), Inches(12.3), Inches(0.3),
              size=Pt(8), color=MUTED, wrap=False)


# ── Public entry point ──────────────────────────────────────────────────────────

def to_pptx(digest: Digest, output_path: str) -> str:
    """
    Generate a branded .pptx digest file from a Digest object.

    Args:
        digest:      Digest object from the pipeline
        output_path: Full path to write the .pptx file

    Returns:
        output_path (for chaining / logging)
    """
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    _build_cover(prs, digest)

    total = len(digest.top_entries)
    for i, entry in enumerate(digest.top_entries, 1):
        _build_signal_slide(prs, entry, i, total)

    prs.save(output_path)
    print(f"[PPTX] Saved: {output_path}")
    return output_path
